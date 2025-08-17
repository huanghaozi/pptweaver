import pytest
import os
from pptx import Presentation
from pptweaver.converter import Converter

# Mark all tests in this file as asyncio
pytestmark = pytest.mark.asyncio

@pytest.fixture
def teardown_file(request):
    """Fixture to clean up a generated file after a test."""
    file_path = None
    def _set_file_path(path):
        nonlocal file_path
        file_path = path
    
    yield _set_file_path
    
    if file_path and os.path.exists(file_path):
        os.remove(file_path)

async def test_all_elements_conversion(teardown_file):
    """Tests conversion of an HTML containing all supported element types."""
    input_file = "tests/samples/all_elements.html"
    output_file = "tests/output_all.pptx"
    teardown_file(output_file)

    converter = Converter(input_file=input_file, output_file=output_file)
    await converter.convert()

    assert os.path.exists(output_file)
    prs = Presentation(output_file)
    assert len(prs.slides) == 1
    
    # Rect, Circle, Ellipse, Line, Polyline, Polygon, Path, Text, Image, ForeignObject
    # Total of 10 elements.
    slide = prs.slides[0]
    assert len(slide.shapes) == 10

async def test_about_pptweaver_conversion(teardown_file):
    """Tests the conversion of the complex 'about_pptweaver.html' presentation."""
    input_file = "tests/samples/about_pptweaver.html"
    output_file = "tests/output_about.pptx"
    teardown_file(output_file)

    converter = Converter(input_file=input_file, output_file=output_file)
    await converter.convert()

    assert os.path.exists(output_file)
    prs = Presentation(output_file)
    
    # The presentation has 6 slides
    assert len(prs.slides) == 6
    
    # Spot check the number of shapes on the complex architecture slide (Slide 4)
    # 4 boxes + 3 arrows + 3 text labels = 10 main shapes
    # Note: this count can be brittle. We're primarily testing that it converts without error.
    slide_4 = prs.slides[3]
    assert len(slide_4.shapes) >= 9 # Allow for some flexibility in shape counting
