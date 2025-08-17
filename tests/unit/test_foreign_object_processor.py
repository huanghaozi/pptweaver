import pytest
from unittest.mock import Mock, MagicMock, call
from pptx.util import Emu, Pt
from pptweaver.elements.foreign_object_processor import ForeignObjectProcessor

@pytest.fixture
def mock_converter():
    converter = Mock()
    converter.px_to_emu.side_effect = lambda px: Emu(int(px * 12700))
    converter.px_to_pt.side_effect = lambda px: Pt(px * 0.75)
    converter.parse_color.side_effect = lambda color_str: ((0, 0, 0), 1.0) # Default mock color
    converter.scale_x = 1.0
    converter.scale_y = 1.0
    return converter

@pytest.fixture
def mock_slide_for_foreign_object():
    slide = Mock()
    mock_shape = Mock()
    mock_text_frame = MagicMock()
    
    # Each time add_paragraph is called, it will return a new MagicMock instance
    # This allows us to track calls on each created paragraph.
    created_paragraphs = []
    def add_paragraph_side_effect():
        p = MagicMock()
        created_paragraphs.append(p)
        return p
    mock_text_frame.add_paragraph.side_effect = add_paragraph_side_effect

    mock_shape.text_frame = mock_text_frame
    slide.shapes.add_textbox.return_value = mock_shape
    return slide, mock_text_frame, created_paragraphs

def test_foreign_object_processor_rich_text(mock_converter, mock_slide_for_foreign_object):
    """Tests processing of a <foreignObject> with rich text and a list."""
    slide_mock, text_frame_mock, created_paragraphs = mock_slide_for_foreign_object
    
    element_data = {
        'rect': {'x': 10, 'y': 20, 'width': 500, 'height': 300},
        'style': {'fontSize': '18px'},
        'text_runs': [
            {'text': 'This is a list:', 'tagName': 'p', 'style': {}},
            {'text': 'First item', 'tagName': 'li', 'style': {}},
            {'text': 'Second item', 'tagName': 'li', 'style': {'fontWeight': 'bold'}},
        ]
    }

    processor = ForeignObjectProcessor(element_data, slide_mock, mock_converter)
    processor.process()

    # 1. Verify a single textbox was created with correct dimensions
    slide_mock.shapes.add_textbox.assert_called_once_with(
        Emu(10 * 12700), Emu(20 * 12700), Emu(500 * 12700), Emu(300 * 12700)
    )

    # 2. Verify the text frame was cleared
    text_frame_mock.clear.assert_called_once()

    # Verify paragraphs were added
    assert text_frame_mock.add_paragraph.call_count == 3
    assert len(created_paragraphs) == 3
    p1_mock, p2_mock, p3_mock = created_paragraphs

    # Verify the runs on the first paragraph
    assert p1_mock.add_run.call_count == 1
    p1_mock.add_run.return_value.text = 'This is a list:'
    
    # Verify the runs on the second paragraph (first list item)
    assert p2_mock.add_run.call_count == 2
    p2_mock.add_run.call_args_list[0].return_value.text = '• '
    p2_mock.add_run.call_args_list[1].return_value.text = 'First item'

    # Verify the runs on the third paragraph (second list item)
    assert p3_mock.add_run.call_count == 2
    p3_mock.add_run.call_args_list[0].return_value.text = '• '
    p3_mock.add_run.call_args_list[1].return_value.text = 'Second item'
