import pytest
import requests_mock
from unittest.mock import Mock
from pptx.util import Emu
from io import BytesIO
import base64
from pptweaver.elements.image_processor import ImageProcessor

# A simple 1x1 transparent PNG as a base64 string
TINY_PNG_BASE64 = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII="

@pytest.fixture
def mock_converter():
    converter = Mock()
    converter.px_to_emu.side_effect = lambda px: Emu(int(px * 12700))
    converter.scale_x = 1.0
    converter.scale_y = 1.0
    return converter

def test_image_processor_data_uri(mock_converter):
    """Tests processing an <image> with a base64 data URI."""
    mock_slide = Mock()
    element_data = {
        'rect': {'x': 10, 'y': 20, 'width': 100, 'height': 80},
        'attributes': {'href': f"data:image/png;base64,{TINY_PNG_BASE64}"}
    }
    
    # Decode the base64 string here to ensure consistency
    expected_bytes = base64.b64decode(TINY_PNG_BASE64)

    processor = ImageProcessor(element_data, mock_slide, mock_converter)
    processor.process()

    mock_slide.shapes.add_picture.assert_called_once()
    args, _ = mock_slide.shapes.add_picture.call_args
    
    # Check that the first arg is a file-like object with the correct bytes
    image_stream = args[0]
    assert isinstance(image_stream, BytesIO)
    assert image_stream.read() == expected_bytes
    
    # Check position and size
    assert args[1] == Emu(10 * 12700) # left
    assert args[2] == Emu(20 * 12700) # top
    assert args[3] == Emu(100 * 12700) # width
    assert args[4] == Emu(80 * 12700) # height

def test_image_processor_url(mock_converter):
    """Tests processing an <image> with an external URL."""
    mock_slide = Mock()
    image_url = "http://example.com/image.png"
    element_data = {
        'rect': {'x': 10, 'y': 20, 'width': 100, 'height': 80},
        'attributes': {'href': image_url}
    }
    
    # Decode the base64 string here to ensure consistency
    expected_bytes = base64.b64decode(TINY_PNG_BASE64)

    with requests_mock.Mocker() as m:
        m.get(image_url, content=expected_bytes)
        
        processor = ImageProcessor(element_data, mock_slide, mock_converter)
        processor.process()

    mock_slide.shapes.add_picture.assert_called_once()
    args, _ = mock_slide.shapes.add_picture.call_args
    
    image_stream = args[0]
    assert isinstance(image_stream, BytesIO)
    assert image_stream.read() == expected_bytes
