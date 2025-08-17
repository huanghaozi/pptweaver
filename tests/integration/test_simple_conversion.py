import pytest
import os
from pptx import Presentation
from pptweaver.converter import Converter

# Mark all tests in this file as asyncio
pytestmark = pytest.mark.asyncio

@pytest.fixture
def setup_teardown_files():
    """Fixture to manage test files, creating them before a test and cleaning up after."""
    input_file = "tests/samples/simple.html"
    output_file = "tests/output_simple.pptx"
    
    # Ensure the output file does not exist before the test
    if os.path.exists(output_file):
        os.remove(output_file)
        
    yield input_file, output_file
    
    # Teardown: remove the output file after the test
    if os.path.exists(output_file):
        os.remove(output_file)

async def test_simple_html_to_pptx_conversion(setup_teardown_files):
    """
    Tests the end-to-end conversion of a simple HTML file to a PPTX presentation.
    """
    input_html, output_pptx = setup_teardown_files

    # 1. Run the conversion
    converter = Converter(input_file=input_html, output_file=output_pptx)
    await converter.convert()

    # 2. Validate the output file
    # Check if the output file was created
    assert os.path.exists(output_pptx), "The output PPTX file was not created."

    # Open the presentation and perform checks
    prs = Presentation(output_pptx)

    # Check the number of slides
    assert len(prs.slides) == 1, "The presentation should have exactly one slide."

    # Check the number of shapes on the first slide
    slide = prs.slides[0]
    # The simple.html file contains a <rect> and a <text> element.
    assert len(slide.shapes) == 2, "The slide should contain exactly two shapes."

    # Optional: More detailed checks on the shapes can be added here
    # For example, checking the text of the textbox or the color of the rectangle.
