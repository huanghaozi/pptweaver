import pytest
from unittest.mock import Mock, MagicMock
from pptx.util import Emu
from pptweaver.elements.line_processor import LineProcessor

@pytest.fixture
def mock_converter():
    """Fixture to create a mock converter with utility methods."""
    converter = Mock()
    converter.px_to_emu.side_effect = lambda px: Emu(px * 12700)
    # Provide a default return value for parse_color
    converter.parse_color.return_value = ((0, 0, 0), 1.0)
    converter.scale_x = 1.0
    converter.scale_y = 1.0
    return converter

@pytest.fixture
def mock_slide_for_connector():
    """Fixture to create a mock slide that can handle connectors."""
    slide = Mock()
    mock_connector = MagicMock()
    # Mock the underlying XML element for arrowhead testing
    mock_connector.line._get_or_add_ln.return_value = MagicMock()
    slide.shapes.add_connector.return_value = mock_connector
    return slide, mock_connector

def test_line_processor_basic_line(mock_converter, mock_slide_for_connector):
    """Tests the processing of a simple <line> element."""
    slide_mock, connector_mock = mock_slide_for_connector
    element_data = {
        'attributes': {'x1': '10', 'y1': '20', 'x2': '100', 'y2': '120'},
        'style': {}
    }

    processor = LineProcessor(element_data, slide_mock, mock_converter)
    processor.process()

    slide_mock.shapes.add_connector.assert_called_once()
    args, _ = slide_mock.shapes.add_connector.call_args
    # args[0] is the connector type, args[1-4] are the coordinates
    assert args[1] == Emu(10 * 12700)
    assert args[2] == Emu(20 * 12700)
    assert args[3] == Emu(100 * 12700)
    assert args[4] == Emu(120 * 12700)
    
    # Ensure arrowhead logic was NOT triggered
    connector_mock.line._get_or_add_ln.assert_not_called()

def test_line_processor_with_arrowhead(mock_converter, mock_slide_for_connector):
    """Tests the processing of a <line> element with a marker-end style."""
    slide_mock, connector_mock = mock_slide_for_connector
    element_data = {
        'attributes': {'x1': '10', 'y1': '20', 'x2': '100', 'y2': '120'},
        'style': {'markerEnd': 'url(#arrow)'} # The value just needs to be non-empty
    }

    processor = LineProcessor(element_data, slide_mock, mock_converter)
    processor.process()

    slide_mock.shapes.add_connector.assert_called_once()
    
    # Ensure arrowhead logic WAS triggered
    connector_mock.line._get_or_add_ln.assert_called_once()
    # Check if the XML append was called
    ln_element = connector_mock.line._get_or_add_ln.return_value
    ln_element.append.assert_called_once()
