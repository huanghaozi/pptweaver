import pytest
from unittest.mock import Mock
from pptx.util import Emu
from pptweaver.elements.oval_processor import OvalProcessor

@pytest.fixture
def mock_converter():
    """Fixture to create a mock converter with utility methods."""
    converter = Mock()
    converter.px_to_emu.side_effect = lambda px: Emu(px * 12700)
    # Provide a default return value for parse_color
    converter.parse_color.return_value = ((0, 0, 255), 1.0)
    converter.scale_x = 1.0
    converter.scale_y = 1.0
    return converter

def test_oval_processor_parsing(mock_converter):
    """
    Tests if the OvalProcessor correctly processes element data for an ellipse/circle.
    """
    element_data = {
        # SVG defines ellipses by a center (cx, cy) and radii (rx, ry).
        # However, the bounding box ('rect') is what our BaseProcessor uses.
        'rect': {'x': 50, 'y': 75, 'width': 200, 'height': 100},
        'style': {'fill': 'rgb(0, 0, 255)'},
        'attributes': {}
    }
    
    mock_slide = Mock()

    # The process method is simple enough that we can infer its behavior
    # by checking the arguments passed to add_shape.
    processor = OvalProcessor(element_data, mock_slide, mock_converter)
    processor.process()

    # Verify that add_shape was called with the correct parameters
    mock_slide.shapes.add_shape.assert_called_once()
    args, _ = mock_slide.shapes.add_shape.call_args
    
    # args[0] is the shape type
    from pptx.enum.shapes import MSO_SHAPE
    assert args[0] == MSO_SHAPE.OVAL
    
    # args[1-4] are the coordinates from the bounding box
    assert args[1] == Emu(50 * 12700)   # left
    assert args[2] == Emu(75 * 12700)   # top
    assert args[3] == Emu(200 * 12700)  # width
    assert args[4] == Emu(100 * 12700)  # height
