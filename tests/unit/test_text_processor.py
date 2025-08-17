import pytest
from unittest.mock import Mock, MagicMock
from pptx.util import Emu, Pt
from pptweaver.elements.text_processor import TextProcessor

@pytest.fixture
def mock_converter():
    """Fixture to create a mock converter with utility methods."""
    converter = Mock()
    converter.px_to_emu.side_effect = lambda px: Emu(px * 12700)
    converter.px_to_pt.side_effect = lambda px: Pt(px * 0.75)
    converter.parse_color.side_effect = lambda color_str: ((0, 128, 0), 1.0)
    converter.scale_x = 1.0
    converter.scale_y = 1.0
    return converter

@pytest.fixture
def mock_slide_for_text():
    """A more detailed mock for a slide, specific to the TextProcessor's needs."""
    slide = Mock()
    mock_shape = Mock()
    
    # Configure the mock text_frame and its paragraphs
    mock_text_frame = MagicMock()
    mock_paragraph = MagicMock()
    
    # Make paragraphs subscriptable and have it return another mock
    mock_text_frame.paragraphs.__getitem__.return_value = mock_paragraph
    
    mock_shape.text_frame = mock_text_frame
    slide.shapes.add_textbox.return_value = mock_shape
    # Return both the slide and the paragraph mock to be used in the test
    return slide, mock_paragraph

def test_text_processor_attribute_parsing(mock_converter, mock_slide_for_text):
    """
    Tests if the TextProcessor correctly parses attributes from the element data.
    """
    # Unpack the fixture tuple
    slide_mock, paragraph_mock = mock_slide_for_text

    element_data = {
        'text': 'Hello PPT',
        'rect': {'x': 100, 'y': 50, 'width': 300, 'height': 50},
        'style': {'fill': 'rgb(0, 128, 0)', 'fontSize': '24px', 'fontWeight': 'bold'},
        'attributes': {}
    }
    
    processor = TextProcessor(element_data, slide_mock, mock_converter)
    processor.process()

    assert processor.text == 'Hello PPT'
    assert processor.left == Emu(100 * 12700)
    assert processor.top == Emu(50 * 12700)
    
    # Verify the mock was used as expected
    slide_mock.shapes.add_textbox.assert_called_once()
    assert paragraph_mock.add_run.called
