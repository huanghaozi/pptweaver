import pytest
from unittest.mock import Mock, MagicMock
from pptx.util import Emu
from pptweaver.elements.polygon_processor import PolygonProcessor

@pytest.fixture
def mock_converter():
    """Fixture to create a mock converter with utility methods."""
    converter = Mock()
    converter.px_to_emu.side_effect = lambda px: Emu(int(px * 12700))
    converter.parse_color.return_value = ((0,0,0), 1.0) # Add default return value
    converter.scale_x = 1.0
    converter.scale_y = 1.0
    return converter

@pytest.fixture
def mock_slide_for_freeform():
    """Fixture to create a mock slide that can handle freeform shapes."""
    slide = Mock()
    freeform_builder = MagicMock()
    slide.shapes.build_freeform.return_value = freeform_builder
    return slide, freeform_builder

def test_polygon_processor_polyline(mock_converter, mock_slide_for_freeform):
    """Tests processing a <polyline> (open shape)."""
    slide_mock, freeform_builder_mock = mock_slide_for_freeform
    element_data = {
        'tagName': 'polyline',
        'attributes': {'points': '10,20 50,60 100,20'},
        'style': {}
    }

    processor = PolygonProcessor(element_data, slide_mock, mock_converter)
    processor.process()

    # Verify build_freeform was started correctly
    slide_mock.shapes.build_freeform.assert_called_once_with(
        Emu(10 * 12700), Emu(20 * 12700)
    )
    # Verify add_line_segments was called with the rest of the points
    freeform_builder_mock.add_line_segments.assert_called_once_with([
        (Emu(50 * 12700), Emu(60 * 12700)),
        (Emu(100 * 12700), Emu(20 * 12700))
    ])
    # Verify the shape was finalized
    freeform_builder_mock.convert_to_shape.assert_called_once()

def test_polygon_processor_polygon(mock_converter, mock_slide_for_freeform):
    """Tests processing a <polygon> (closed shape)."""
    slide_mock, freeform_builder_mock = mock_slide_for_freeform
    element_data = {
        'tagName': 'polygon',
        'attributes': {'points': '10,20 50,60 100,20'},
        'style': {}
    }

    processor = PolygonProcessor(element_data, slide_mock, mock_converter)
    processor.process()

    # For a polygon, the starting point must be appended to the list to close it
    expected_point_list = [
        (Emu(50 * 12700), Emu(60 * 12700)),
        (Emu(100 * 12700), Emu(20 * 12700)),
        (Emu(10 * 12700), Emu(20 * 12700)) # Closing point
    ]

    # Verify build_freeform was started correctly
    slide_mock.shapes.build_freeform.assert_called_once_with(
        Emu(10 * 12700), Emu(20 * 12700)
    )
    # Verify add_line_segments was called with the rest of the points
    freeform_builder_mock.add_line_segments.assert_called_once_with(expected_point_list)
    # Verify the shape was finalized
    freeform_builder_mock.convert_to_shape.assert_called_once()
